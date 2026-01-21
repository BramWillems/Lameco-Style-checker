# layout_rules.py

from docx import Document
from pptx import Presentation

EMU_PER_CM = 360000

# ============================================================
# POSITION CLASSIFICATION
# ============================================================

def classify_position(x, y, width, height):
    horizontal = (
        "left" if x < width / 3 else
        "center" if x < 2 * width / 3 else
        "right"
    )

    vertical = (
        "top" if y < height / 3 else
        "middle" if y < 2 * height / 3 else
        "bottom"
    )

    return f"{vertical}-{horizontal}"

# ============================================================
# NORMALIZED EXTRACTORS
# ============================================================

def extract_logos(path):
    if path.endswith(".pptx"):
        return _extract_logos_pptx(path)
    elif path.endswith(".docx"):
        return _extract_logos_docx(path)
    else:
        raise ValueError("Unsupported file type")

# ---------------- PPTX ----------------

def _extract_logos_pptx(path):
    prs = Presentation(path)
    logos = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                logos.append({
                    "page": slide_index,
                    "x": shape.left,
                    "y": shape.top,
                    "page_width": prs.slide_width,
                    "page_height": prs.slide_height,
                    "source": "pptx"
                })

    return logos

# ---------------- DOCX ----------------

def _extract_logos_docx(path):
    doc = Document(path)
    logos = []

    # Page size (approximate, first section)
    section = doc.sections[0]
    page_width = section.page_width
    page_height = section.page_height

    # Floating images (anchors)
    anchors = doc.part._element.xpath('//wp:anchor')
    for anchor in anchors:
        try:
            x = int(anchor.xpath('.//wp:positionH/wp:posOffset')[0].text)
            y = int(anchor.xpath('.//wp:positionV/wp:posOffset')[0].text)

            logos.append({
                "page": 1,
                "x": x,
                "y": y,
                "page_width": page_width,
                "page_height": page_height,
                "source": "docx"
            })
        except Exception:
            continue

    return logos

# ============================================================
# RULE: LOGO POSITION
# ============================================================

def check_logo_position(path, expected_position):
    logos = extract_logos(path)
    errors = []

    if not logos:
        errors.append({
            "type": "logo_missing",
            "severity": "high",
            "message": "No logo found in document"
        })
        return errors

    for logo in logos:
        actual_position = classify_position(
            logo["x"],
            logo["y"],
            logo["page_width"],
            logo["page_height"]
        )

        if actual_position != expected_position:
            errors.append({
                "type": "logo_position",
                "severity": "medium",
                "page": logo["page"],
                "expected": expected_position,
                "found": actual_position,
                "source": logo["source"]
            })

    return errors
