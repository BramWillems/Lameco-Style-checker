# app/services/analyzer.py

import os
import re
import requests
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple

from docx import Document
from pptx import Presentation
from requests.exceptions import HTTPError, RequestException
from transformers import pipeline

# ==============================================================================
# CONFIG (hard-coded for demo)
# ==============================================================================
SAPLING_API_KEY = "TJ82T8FXJU7H5LIQ8FII9WKBIAFN9OWG"
SAPLING_ENDPOINT = "https://api.sapling.ai/api/v1/edits"
TARGET_FONT = "Calibri"

# ==============================================================================
# TONE OF VOICE – AI (lazy loaded)
# ==============================================================================
_tone_classifier = None


def get_tone_classifier():
    global _tone_classifier
    if _tone_classifier is None:
        _tone_classifier = pipeline(
            "text-classification",
            model="cardiffnlp/twitter-roberta-base-sentiment-latest"
        )
    return _tone_classifier


# In analyzer.py, VERANDER DEZE FUNCTIE:
def check_tone_of_voice_ai(text: str) -> str:
    """
    AI-based Tone of Voice detection.
    Returns: 'formal', 'neutral', or 'informal'
    """
    if not text or len(text.strip()) < 10:
        return "neutral"
    
    text_lower = text.lower()
    
    # Formele indicatoren
    formal_words = ["geachte", "verzoeken", "vriendelijk", "hoogachtend", 
                   "bij dezen", "ondergetekende", "vertrouwen", "achtzaam"]
    
    # Informele indicatoren
    informal_words = ["hey", "hoi", "ff", "even", "doei", "oké", "ok", 
                     "graag gedaan", "checken", "makkelijk"]
    
    # Controleer op formele taal
    formal_count = sum(1 for word in formal_words if word in text_lower)
    informal_count = sum(1 for word in informal_words if word in text_lower)
    
    # Controleer op uitroeptekens (informeel)
    exclamation_count = text.count('!')
    
    # Bepaal toon
    if informal_count > 0 or exclamation_count > 1:
        return "informal"
    elif formal_count > 0:
        return "formal"
    else:
        # Gebruik nog steeds de AI classifier als backup
        try:
            classifier = get_tone_classifier()
            result = classifier(text[:512])[0]
            label = result["label"].lower()
            
            if label in ["negative", "informal"]:
                return "informal"
            elif label == "positive":
                return "formal"
        except:
            pass
        
        return "neutral"


# ==============================================================================
# DOCX FONT UTILITIES
# ==============================================================================
def get_all_runs(doc: Document) -> List[Dict[str, Any]]:
    runs: List[Dict[str, Any]] = []
    para_counter = 0
    current_section: Optional[str] = None

    def collect_paragraphs(paragraphs):
        nonlocal para_counter, current_section
        for para in paragraphs:
            para_counter += 1
            style_name = para.style.name if para.style else ""

            if style_name.startswith("Heading"):
                current_section = para.text.strip()
            elif current_section is None and para.text.strip():
                current_section = para.text.strip()

            for run in para.runs:
                if run.text and run.text.strip():
                    runs.append({
                        "paragraph": para_counter,
                        "section": current_section,
                        "run": run
                    })

    def collect_tables(tables):
        for table in tables:
            for row in table.rows:
                for cell in row.cells:
                    collect_paragraphs(cell.paragraphs)
                    collect_tables(cell.tables)

    collect_paragraphs(doc.paragraphs)
    collect_tables(doc.tables)
    return runs


def get_font_from_run(run) -> Optional[str]:
    if run.font and run.font.name:
        return run.font.name

    r = run._element
    rPr = r.find(
        ".//w:rPr",
        namespaces={"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
    )
    if rPr is not None:
        rFonts = rPr.find(
            "w:rFonts",
            namespaces={"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
        )
        if rFonts is not None:
            return rFonts.attrib.get(
                "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}ascii"
            )
    return None


def check_font_consistency_docx(path: str, target_font: str) -> List[Dict[str, Any]]:
    doc = Document(path)
    mismatches: List[Dict[str, Any]] = []

    for entry in get_all_runs(doc):
        run = entry["run"]
        font_name = get_font_from_run(run)

        if font_name and font_name != target_font:
            mismatches.append({
                "paragraph": entry["paragraph"],
                "section": entry["section"],
                "text": run.text.strip(),
                "font_found": font_name,
                "font_expected": target_font
            })

    return mismatches


# ==============================================================================
# GRAMMAR CHECK (SAPLING)
# ==============================================================================
# In analyzer.py, VERANDER DEZE FUNCTIE:
def check_text_with_sapling(text: str, lang: str = "nl") -> List[Dict[str, Any]]:
    if not text.strip():
        return []
    
    payload = {
        "key": SAPLING_API_KEY,
        "text": text,
        "lang": lang
    }
    
    try:
        resp = requests.post(SAPLING_ENDPOINT, json=payload, timeout=10)
        
        if resp.status_code == 200:
            return resp.json().get("edits", [])
        else:
            # Fallback: gebruik eenvoudige spellingscontrole
            return get_basic_spelling_check(text, lang)
            
    except Exception as e:
        # Fallback naar basic checking als API faalt
        return get_basic_spelling_check(text, lang)

def get_basic_spelling_check(text: str, lang: str = "nl") -> List[Dict[str, Any]]:
    """Eenvoudige fallback spellingcontrole"""
    common_errors = {
        "ff": "even",
        "mss": "misschien",
        "idd": "inderdaad",
        "iig": "in ieder geval",
        "wss": "wees",
        "wdt": "wordt",
    }
    
    edits = []
    words = text.split()
    
    for i, word in enumerate(words):
        word_lower = word.lower().strip('.,!?;:')
        
        if word_lower in common_errors:
            edits.append({
                "start": text.find(word),
                "end": text.find(word) + len(word),
                "type": "spelling",
                "replacement": common_errors[word_lower],
                "message": f"Mogelijke informele afkorting: '{word}' → '{common_errors[word_lower]}'"
            })
    
    return edits

# ==============================================================================
# GRAMMAR FILTERING
# ==============================================================================
def should_check_grammar(text: str) -> bool:
    text = text.strip()
    if not text:
        return False

    words = text.split()

    if len(words) < 3:
        return False

    if all(not any(c.isalpha() for c in w) for w in words):
        return False

    if "." not in text and "?" not in text and "!" not in text and len(text) < 40:
        return False

    return True


# ==============================================================================
# CONTEXT FORMATTER
# ==============================================================================
def format_error_context(text: str, start: int, end: int) -> Tuple[str, str]:
    start = max(0, start)
    end = min(len(text), end)

    prefix = text[:start]
    error = text[start:end] if end > start else text[start:start + 1]
    suffix = text[end:]

    prefix = ("..." + prefix[-40:]) if len(prefix) > 40 else prefix
    suffix = (suffix[:40] + "...") if len(suffix) > 40 else suffix

    return f"{prefix} **{error}** {suffix}".replace("\n", " "), error


# ==============================================================================
# DOCUMENT READERS
# ==============================================================================
def extract_docx_paragraphs(path: str) -> List[str]:
    doc = Document(path)
    return [p.text.strip() for p in doc.paragraphs if p.text.strip()]


def extract_pptx_blocks(path: str) -> List[Dict[str, Any]]:
    prs = Presentation(path)
    blocks = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and text.strip():
                blocks.append({
                    "slide": slide_idx,
                    "text": text.strip()
                })
    return blocks


# ADD THESE FUNCTIONS TO analyzer.py (somewhere after get_font_from_run)

# ADD THESE FUNCTIONS TO analyzer.py (somewhere after get_font_from_run)

def get_font_weight_from_run(run) -> Optional[str]:
    """Get font weight (bold, italic, etc.) from a run"""
    if not run:
        return None
    
    weight = "regular"
    
    # Check bold
    if run.font and run.font.bold is not None:
        if run.font.bold:
            weight = "bold"
    
    # Check italic
    if run.font and run.font.italic is not None:
        if run.font.italic:
            weight = "italic" if weight == "regular" else f"bold-italic"
    
    # Fallback to XML parsing
    if weight == "regular":
        try:
            r = run._element
            rPr = r.find(
                ".//w:rPr",
                namespaces={"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
            )
            if rPr is not None:
                # Check for bold
                b = rPr.find(
                    "w:b",
                    namespaces={"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
                )
                if b is not None:
                    weight = "bold"
                
                # Check for italic
                i = rPr.find(
                    "w:i",
                    namespaces={"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
                )
                if i is not None:
                    weight = "italic" if weight == "regular" else f"bold-italic"
        except:
            pass
    
    return weight

def check_font_weight_consistency_docx(path: str, expected_weight: str = "regular") -> List[Dict[str, Any]]:
    """Check font weight consistency in DOCX"""
    from docx import Document
    doc = Document(path)
    mismatches = []
    
    # Get all runs with the function from your existing code
    for entry in get_all_runs(doc):  # Assuming get_all_runs exists
        run = entry["run"]
        actual_weight = get_font_weight_from_run(run)
        
        if actual_weight and actual_weight != expected_weight:
            # For headings, bold might be acceptable
            if "heading" in str(entry.get("section", "")).lower() and actual_weight == "bold":
                continue
                
            mismatches.append({
                "paragraph": entry["paragraph"],
                "section": entry["section"],
                "text": run.text.strip()[:100],
                "weight_found": actual_weight,
                "weight_expected": expected_weight,
                "context": entry.get("section", "")
            })
    
    return mismatches

# Then update the check_font_consistency_docx function to include weight checking:
def check_font_consistency_docx(path: str, target_font: str, check_weight: bool = True) -> List[Dict[str, Any]]:
    """Check font family and weight consistency in DOCX"""
    doc = Document(path)
    mismatches: List[Dict[str, Any]] = []
    
    for entry in get_all_runs(doc):
        run = entry["run"]
        font_name = get_font_from_run(run)
        
        # Check font family
        if font_name and font_name != target_font:
            mismatches.append({
                "type": "font_family",
                "paragraph": entry["paragraph"],
                "section": entry["section"],
                "text": run.text.strip()[:100],
                "font_found": font_name,
                "font_expected": target_font,
                "severity": "high"
            })
        
        # Check font weight (if enabled)
        if check_weight:
            font_weight = get_font_weight_from_run(run)
            # For body text, expect regular weight
            # Headings can be bold
            style_name = ""
            if hasattr(run, 'parent') and hasattr(run.parent, 'style'):
                style_name = run.parent.style.name if run.parent.style else ""
            
            is_heading = "heading" in style_name.lower() or "title" in style_name.lower()
            
            if font_weight and font_weight == "bold" and not is_heading:
                mismatches.append({
                    "type": "font_weight",
                    "paragraph": entry["paragraph"],
                    "section": entry["section"],
                    "text": run.text.strip()[:100],
                    "weight_found": font_weight,
                    "weight_expected": "regular",
                    "severity": "medium",
                    "context": "Body text should not be bold"
                })
            elif font_weight and font_weight == "italic" and is_heading:
                mismatches.append({
                    "type": "font_weight",
                    "paragraph": entry["paragraph"],
                    "section": entry["section"],
                    "text": run.text.strip()[:100],
                    "weight_found": font_weight,
                    "weight_expected": "bold or regular",
                    "severity": "medium",
                    "context": "Headings should not be italic"
                })
    
    return mismatches

# ==============================================================================
# NORMALIZATION FOR UI & PDF
# ==============================================================================

def normalize_to_ui(report: Dict[str, Any], target_font: str) -> Dict[str, Any]:
    findings: List[Dict[str, Any]] = []
    is_pptx = report["bestand"].lower().endswith(".pptx")
    
    # Counter voor unieke IDs
    issue_counter = 0
    
    # Font issues
    for f in report.get("font_controle", []):
        issue_counter += 1
        issue_type = f.get("type", "font_family")
        if issue_type == "font_family":
            category = "Font Issues"
            message = f"Expected {target_font}, found {f['font_found']}"
        else:  # font_weight
            category = "Font Weight"
            message = f"Expected {f.get('weight_expected', 'regular')}, found {f.get('weight_found')}"
        
        # Simpele ID zonder hash
        element_id = f"issue-{issue_counter}"
        
        findings.append({
            "location": f"Paragraph {f['paragraph']}",
            "page": f["paragraph"],
            "category": category,
            "severity": f.get("severity", "medium"),
            "message": message,
            "detail": f.get("text", "") or f.get("context", ""),
            "hyperlink": f"#{element_id}",  # Directe hash naar ID
            "element_id": element_id
        })
    
    # Color issues
    for c in report.get("color_issues", []):
        issue_counter += 1
        element_id = f"issue-{issue_counter}"
        
        findings.append({
            "location": c.get("location", "Unknown"),
            "page": 1,
            "category": "Color Issues",
            "severity": c.get("severity", "medium"),
            "message": f"Expected brand color {c.get('expected_color')}, found {c.get('color_found')}",
            "detail": f"{c.get('element', '')}: {c.get('context', '')}",
            "hyperlink": f"#{element_id}",
            "element_id": element_id
        })
    
    # Logo issues
    for l in report.get("logo_issues", []):
        issue_counter += 1
        element_id = f"issue-{issue_counter}"
        
        findings.append({
            "location": l.get("location", "Document"),
            "page": l.get("slide", 1) if is_pptx else 1,
            "category": "Logo Issues",
            "severity": l.get("severity", "medium"),
            "message": l.get("message", ""),
            "detail": l.get("details", ""),
            "hyperlink": f"#{element_id}",
            "element_id": element_id
        })
    
    # Grammar issues
    for g in report.get("grammatica_fouten", []):
        issue_counter += 1
        label = "Slide" if is_pptx else "Paragraph"
        location_num = g["alinea"]
        element_id = f"issue-{issue_counter}"
        
        findings.append({
            "location": f"{label} {location_num}",
            "page": location_num,
            "category": "Grammar & Spelling",
            "severity": "low",
            "message": g.get("message") or f"Grammar issue: {g['fout']}",
            "detail": g.get("context", ""),
            "hyperlink": f"#{element_id}",
            "element_id": element_id
        })

    
    # Tone issues (existing)
    for t in report.get("tone_of_voice", []):
        label = "Slide" if is_pptx else "Paragraph"
        location_num = t["alinea"]
        hyperlink = f"#location={label.lower()}-{location_num}"
        
        findings.append({
            "location": f"{label} {location_num}",
            "page": location_num,
            "category": "Tone of Voice",
            "severity": "medium",
            "message": t["message"],
            "detail": "",
            "hyperlink": hyperlink,
            "element_id": f"tone-{location_num}"
        })
    
    # Calculate score and counts
    total = len(findings)
    score = max(0, 100 - total * 2)
    
    counts = {
        "high": sum(f["severity"] == "high" for f in findings),
        "medium": sum(f["severity"] == "medium" for f in findings),
        "low": sum(f["severity"] == "low" for f in findings),
        "total": total
    }
    
    # Update categories
    categories = [
        "Font Issues",
        "Font Weight",
        "Color Issues",
        "Logo Issues",
        "Spacing",
        "Grammar & Spelling",
        "Tone of Voice"
    ]
    
    category_counts = {
        c: {"high": 0, "medium": 0, "low": 0} for c in categories
    }
    
    for f in findings:
        cat = f["category"]
        if cat in category_counts:
            category_counts[cat][f["severity"]] += 1
    
    return {
        "document": report["bestand"],
        "created_at": report["datum"],
        "score": score,
        "counts": counts,
        "category_counts": category_counts,
        "findings": findings,
        "raw": report,
        "has_hyperlinks": True
    }

# ==============================================================================
# MAIN PIPELINE
# ==============================================================================
def analyze_document(
    path: str,
    target_font: str = TARGET_FONT,
    check_fonts: bool = True,
    check_spelling: bool = True,
    check_tone: bool = True,
    check_colors: bool = True,  # NEW
    check_logo: bool = True,    # NEW
    lang: str = "nl"
) -> Dict[str, Any]:

    ext = os.path.splitext(path)[1].lower()

    report = {
        "bestand": os.path.basename(path),
        "datum": datetime.now().isoformat(),
        "font_controle": [],
        "grammatica_fouten": [],
        "tone_of_voice": [], 
        "color_issues": [],  # NEW
        "logo_issues": []    # NEW
    }

    if ext == ".docx":
        paragraphs = extract_docx_paragraphs(path)

        if check_fonts:
            report["font_controle"] = check_font_consistency_docx(path, target_font, check_weight=True)


        for idx, text in enumerate(paragraphs, start=1):
            if check_spelling and should_check_grammar(text):
                for issue in check_text_with_sapling(text, lang):
                    context, fout = format_error_context(
                        text, issue.get("start", 0), issue.get("end", 0)
                    )
                    report["grammatica_fouten"].append({
                        "alinea": idx,
                        "fout": fout,
                        "suggestie": issue.get("replacement"),
                        "type": issue.get("type"),
                        "context": context,
                        "message": issue.get("message")
                    })

            if check_tone:
                tone = check_tone_of_voice_ai(text)
                if tone != "neutral":
                    report["tone_of_voice"].append({
                        "alinea": idx,
                        "tone": tone,
                        "message": f"Tone of voice detected as '{tone}'"
                    })

    elif ext == ".pptx":
        for block in extract_pptx_blocks(path):
            text = block["text"]

            if check_spelling and should_check_grammar(text):
                for issue in check_text_with_sapling(text, lang):
                    context, fout = format_error_context(
                        text, issue.get("start", 0), issue.get("end", 0)
                    )
                    report["grammatica_fouten"].append({
                        "alinea": block["slide"],
                        "fout": fout,
                        "suggestie": issue.get("replacement"),
                        "type": issue.get("type"),
                        "context": context,
                        "message": issue.get("message")
                    })

            if check_tone:
                tone = check_tone_of_voice_ai(text)
                if tone != "neutral":
                    report["tone_of_voice"].append({
                        "alinea": block["slide"],
                        "tone": tone,
                        "message": f"Tone of voice detected as '{tone}'"
                    })

    else:
        raise ValueError("Only .docx and .pptx supported")

    report["totaal_grammatica_fouten"] = len(report["grammatica_fouten"])
    return report
