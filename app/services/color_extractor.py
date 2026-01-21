# app/services/color_extractor.py
import re
from typing import List, Dict, Any, Optional, Tuple
from docx import Document
from pptx import Presentation
from pptx.dml.color import RGBColor

def rgb_to_hex(r: int, g: int, b: int) -> str:
    """Convert RGB to hex color code"""
    return f"#{r:02x}{g:02x}{b:02x}".upper()

def hex_to_rgb(hex_color: str) -> Optional[Tuple[int, int, int]]:
    """Convert hex color to RGB"""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 6:
        try:
            return (
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
        except ValueError:
            return None
    return None

def is_color_similar(color1: str, color2: str, threshold: int = 20) -> bool:
    """Check if two colors are similar within threshold"""
    rgb1 = hex_to_rgb(color1)
    rgb2 = hex_to_rgb(color2)
    
    if not rgb1 or not rgb2:
        return False
    
    # Calculate Euclidean distance in RGB space
    distance = sum((c1 - c2) ** 2 for c1, c2 in zip(rgb1, rgb2)) ** 0.5
    return distance <= threshold

def extract_colors_docx(path: str) -> List[Dict[str, Any]]:
    """Extract colors from DOCX document - SIMPLIFIED VERSION"""
    colors = []
    try:
        doc = Document(path)
        
        # Check paragraph colors
        for para_idx, paragraph in enumerate(doc.paragraphs, 1):
            # Text color
            for run in paragraph.runs:
                if run.font and run.font.color and run.font.color.rgb:
                    try:
                        # Get RGB values
                        rgb = run.font.color.rgb
                        if isinstance(rgb, tuple) and len(rgb) == 3:
                            hex_color = rgb_to_hex(rgb[0], rgb[1], rgb[2])
                            colors.append({
                                "type": "text",
                                "location": f"Paragraph {para_idx}",
                                "element": run.text[:30] + "..." if len(run.text) > 30 else run.text,
                                "color_hex": hex_color,
                                "context": paragraph.text[:50] + "..." if len(paragraph.text) > 50 else paragraph.text
                            })
                    except Exception as e:
                        print(f"Color extraction warning: {e}")
                        continue
        
        # Skip table color extraction for now to avoid XML errors
        
    except Exception as e:
        print(f"Color extraction error: {e}")
    
    return colors

def extract_colors_pptx(path: str) -> List[Dict[str, Any]]:
    """Extract colors from PPTX presentation"""
    colors = []
    prs = Presentation(path)
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape_idx, shape in enumerate(slide.shapes, 1):
            # Text color
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font = run.font
                        if font and hasattr(font, 'color'):
                            color = font.color
                            if color and hasattr(color, 'rgb'):
                                if isinstance(color.rgb, RGBColor):
                                    hex_color = rgb_to_hex(color.rgb[0], color.rgb[1], color.rgb[2])
                                    colors.append({
                                        "type": "text",
                                        "location": f"Slide {slide_idx}, Shape {shape_idx}",
                                        "element": run.text[:30] + "..." if len(run.text) > 30 else run.text,
                                        "color_hex": hex_color,
                                        "context": shape.text[:50] + "..." if len(shape.text) > 50 else shape.text
                                    })
            
            # Shape fill color
            if hasattr(shape, 'fill'):
                fill = shape.fill
                if fill and hasattr(fill, 'solid'):
                    if hasattr(fill.solid, 'fore_color'):
                        color = fill.solid.fore_color
                        if hasattr(color, 'rgb'):
                            if isinstance(color.rgb, RGBColor):
                                hex_color = rgb_to_hex(color.rgb[0], color.rgb[1], color.rgb[2])
                                colors.append({
                                    "type": "shape_fill",
                                    "location": f"Slide {slide_idx}, Shape {shape_idx}",
                                    "element": f"{shape.name if hasattr(shape, 'name') else 'Shape'} fill",
                                    "color_hex": hex_color,
                                    "context": ""
                                })
    
    # Check theme colors
    if hasattr(prs, 'slide_master') and hasattr(prs.slide_master, 'color_scheme'):
        scheme = prs.slide_master.color_scheme
        for idx in range(1, 13):  # Standard color scheme has 12 colors
            color = getattr(scheme, f'accent{idx}', None)
            if color and hasattr(color, 'rgb'):
                if isinstance(color.rgb, RGBColor):
                    hex_color = rgb_to_hex(color.rgb[0], color.rgb[1], color.rgb[2])
                    colors.append({
                        "type": "theme",
                        "location": "Presentation Theme",
                        "element": f"Accent {idx}",
                        "color_hex": hex_color,
                        "context": "Theme color"
                    })
    
    return colors

def check_color_compliance(
    extracted_colors: List[Dict[str, Any]], 
    brand_color: str,
    tolerance: int = 20
) -> List[Dict[str, Any]]:
    """Check if colors match brand guidelines"""
    violations = []
    
    for color_info in extracted_colors:
        extracted_hex = color_info.get("color_hex", "")
        if not extracted_hex:
            continue
            
        if not is_color_similar(extracted_hex, brand_color, tolerance):
            violations.append({
                **color_info,
                "expected_color": brand_color,
                "color_found": extracted_hex,
                "severity": "medium" if color_info["type"] == "text" else "low"
            })
    
    return violations