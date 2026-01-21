import requests
from docx import Document
from pptx import Presentation
import os
import json  # Nodig voor de export

# ==============================================================================
# CONFIGURATIE
# ==============================================================================

SAPLING_API_KEY = "TJ82T8FXJU7H5LIQ8FII9WKBIAFN9OWG"
CURRENT_PROVIDER = "sapling" 

ENDPOINTS = {
    "edits": "https://api.sapling.ai/api/v1/edits",
}


FORMAL_INDICATORS = [
    "u", "uw", "geachte", "graag", "vriendelijk", "hoogachtend",
    "wij verzoeken", "wij adviseren", "met betrekking tot"
]

INFORMAL_INDICATORS = [
    "je", "jij", "hey", "hoi", "doei", "thanks", "oké", "ok",
    "we gaan", "laat maar", "ff", "even"
]


# ==============================================================================
# 1. API CLIENT (Grammatica & Tone of Voice)
# ==============================================================================

def check_text_with_sapling(text):
    if not text or not text.strip():
        return {"edits": []}

    url = ENDPOINTS["edits"]
    payload = {
        "key": SAPLING_API_KEY,
        "text": text,
        "lang": "nl"
    }
    response = requests.post(url, json=payload)
    response.raise_for_status()
    return response.json()

def check_text_with_openai(text):
    # Placeholder voor later
    return {"edits": []}

def get_grammar_edits(text):
    if CURRENT_PROVIDER == "sapling":
        return check_text_with_sapling(text)
    elif CURRENT_PROVIDER == "openai":
        return check_text_with_openai(text)
    else:
        raise ValueError("Onbekende provider")

# --- NIEUW: Tone of Voice functie ---

def check_tone_of_voice(text: str) -> str:
    """
    Rule-based Tone of Voice detection.
    Returns: 'formal', 'neutral', or 'informal'
    """
    if not text or len(text.strip()) < 10:
        return "neutral"

    text_lower = text.lower()

    formal_hits = sum(1 for w in FORMAL_INDICATORS if w in text_lower)
    informal_hits = sum(1 for w in INFORMAL_INDICATORS if w in text_lower)

    # Veel uitroeptekens = informeler
    exclamations = text.count("!")

    # Lange zinnen → formeler
    avg_sentence_length = sum(
        len(s.split()) for s in re.split(r"[.!?]", text) if s.strip()
    ) / max(1, len(re.split(r"[.!?]", text)))

    if informal_hits + exclamations > formal_hits:
        return "informal"

    if formal_hits > 0 or avg_sentence_length > 15:
        return "formal"

    return "neutral"




# ==============================================================================
# 2. CONTEXT & OPMAAK
# ==============================================================================

def format_error_context(paragraph_text, start, end):
    """Maakt de fout dikgedrukt binnen de alinea."""
    start = max(0, start)
    end = min(len(paragraph_text), end)

    prefix = paragraph_text[:start]
    error_word = paragraph_text[start:end]
    suffix = paragraph_text[end:]

    short_prefix = prefix[-40:] if len(prefix) > 40 else prefix
    short_suffix = suffix[:40] if len(suffix) > 40 else suffix
    
    disp_prefix = "..." + short_prefix if len(prefix) > 40 else short_prefix
    disp_suffix = short_suffix + "..." if len(suffix) > 40 else short_suffix

    clean_context = f"{disp_prefix} **{error_word}** {disp_suffix}".replace('\n', ' ')
    
    return clean_context, error_word

# ==============================================================================
# 3. DOCUMENT LEZERS
# ==============================================================================

def extract_paragraphs(path):
    if path.endswith(".docx"):
        doc = Document(path)
        return [p.text for p in doc.paragraphs if p.text.strip()]
    
    elif path.endswith(".pptx"):
        prs = Presentation(path)
        text_items = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_items.append(shape.text)
        return text_items
    
    else:
        raise ValueError("Bestand moet .docx of .pptx zijn")

# ==============================================================================
# 4. HOOFD PROGRAMMA
# ==============================================================================

def analyze_document(path):
    print(f"\n📄 Document verwerken: {path}...")
    
    try:
        paragraphs = extract_paragraphs(path)
    except Exception as e:
        print(f"Fout bij lezen bestand: {e}")
        return

    print(f"🔍 Start analyse per alinea ({len(paragraphs)} alinea's)...")
    print("=" * 60)

    # We maken een datastructuur om alles in op te slaan voor de JSON export
    rapport_data = {
        "bestandsnaam": path,
        "analysedatum": "Vandaag", # Kun je later dynamisch maken met datetime
        "tone_of_voice_algemeen": check_tone_of_voice("Hele document tekst hier..."), # De check aanroepen
        "gevonden_fouten": []
    }

    total_errors = 0

    for index, para_text in enumerate(paragraphs, 1):
        try:
            api_result = get_grammar_edits(para_text)
            edits = api_result.get("edits", [])

            if not edits:
                continue

            total_errors += len(edits)

            for issue in edits:
                start = issue['start']
                end = issue['end']
                suggestie = issue['replacement']
                
                context_str, original_error = format_error_context(para_text, start, end)

                if not original_error.strip():
                    continue

                # 1. Printen naar scherm (voor jou)
                print(f"Fout:      '{original_error}'")
                print(f"Suggestie: '{suggestie}'")
                print(f"Locatie:   Alinea {index}")
                print(f"Context:   {context_str}")
                print("-" * 60)

                # 2. Toevoegen aan rapport data (voor de website/JSON)
                fout_object = {
                    "fout": original_error,
                    "suggestie": suggestie,
                    "locatie": f"Alinea {index}",
                    "context": context_str,
                    "technische_info": {
                        "start_index": start,
                        "eind_index": end,
                        "error_type": issue.get("type", "unknown")
                    }
                }
                rapport_data["gevonden_fouten"].append(fout_object)

        except Exception as e:
            print(f"⚠️ Fout bij checken alinea {index}: {e}")

    # Voeg totaal aantal fouten toe aan het rapport
    rapport_data["totaal_aantal_fouten"] = total_errors

    print(f"\nKlaar ✔️ - Totaal {total_errors} suggesties gevonden.")

    # --- NIEUW: JSON Opslaan ---
    json_bestandsnaam = "analyse_rapport.json"
    try:
        with open(json_bestandsnaam, "w", encoding='utf-8') as f:
            json.dump(rapport_data, f, indent=4, ensure_ascii=False)
        print(f"💾 Rapport succesvol opgeslagen als '{json_bestandsnaam}'")
        print("   (Dit bestand kan door de website worden ingelezen)")
        return rapport_data
    except Exception as e:
        print(f"⚠️ Kon JSON rapport niet opslaan: {e}")
        return rapport_data

if __name__ == "__main__":
    errors = analyze_document("Test.docx")
    print(errors)
